from msal import ConfidentialClientApplication
import os
from dotenv import load_dotenv
import requests
import json
from dataclasses import dataclass
from io import BytesIO
from docx import Document
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation
import sqlite3
import chunking
import time

@dataclass
class SharePointFile:
    id: str
    site_id: str
    drive_id: str
    name: str
    mime_type: str
    size: int
    web_url: str
    created_at: str
    modified_at: str
    parent_id: str
    text: str

load_dotenv()

TENANT_ID = os.getenv("TENANT_ID")
CLIENT_ID = os.getenv("CLIENT_ID")
CLIENT_SECRET = os.getenv("CLIENT_SECRET")


AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
SCOPES = ["https://graph.microsoft.com/.default"]

connection = sqlite3.connect("sharpoint.db")

cursor = connection.cursor()

cursor.execute("""
    CREATE TABLE IF NOT EXISTS documents (
        id TEXT PRIMARY KEY,
        site_id TEXT NOT NULL,
        drive_id TEXT NOT NULL,
        name TEXT NOT NULL,
        mime_type TEXT,
        size INTEGER,
        web_url TEXT,
        created_at TEXT,
        modified_at TEXT,
        parent_id TEXT,
        text TEXT
    )
""")

connection.commit()



app = ConfidentialClientApplication(
    CLIENT_ID,
    authority=AUTHORITY,
    client_credential=CLIENT_SECRET
)

def get_headers():
    result = app.acquire_token_for_client(scopes=SCOPES)
    if "access_token" not in result:
        raise Exception("No AccessToken found")
    return {"Authorization": f"Bearer {result['access_token']}"}

def safe_get(url, headers, timeout=30, retries=3):
    for attempt in range(retries):
        try:
            response = requests.get(
                url,
                headers=headers,
                timeout=timeout
            )
            response.raise_for_status()
            return response
        except requests.exceptions.RequestException as e:
            if attempt < retries - 1:
                time.sleep(5)
            else:
                raise

sites_to_process = ["wi3.sharepoint.com,a76fe795-879d-4e8e-ba06-6a0735654277,a59f0326-2de6-4360-8a55-1d0c89783f0e",
                    "wi3.sharepoint.com,9d8c4829-a660-4f44-b273-45fbd13909a0,85dc0cfe-9371-4772-a7b7-87fd1aae1813",
                    "wi3.sharepoint.com,dfc191e5-f9cd-4870-b5c5-4a7f056c1710,a59f0326-2de6-4360-8a55-1d0c89783f0e",
                    "wi3.sharepoint.com,3772b9f6-99b4-44f4-957c-bf1a8afeebf8,b9ce6d80-a819-4435-97b8-e2091452379b",
                    "wi3.sharepoint.com,ee22c8d5-9932-421e-baa5-7d7c7af97e38,a59f0326-2de6-4360-8a55-1d0c89783f0e",
                    "wi3.sharepoint.com,dceed059-1ad1-4c73-922d-aa2866e394ee,4e2acef0-78a2-4297-9b8a-983269dc38d4",
                    "wi3.sharepoint.com,e368bee1-ac0e-428c-b945-4143122ab90c,58836373-e931-475a-8104-bfce36c58b27",
                    "wi3.sharepoint.com,9b5d188c-29bd-4379-aad0-27b1765bdacf,94d9b353-de81-4687-8b3e-db92f4384512",
                    "wi3.sharepoint.com,3aff7cab-2eed-4407-b34d-e0e290ed81ef,ccb943f3-43e5-4d0f-b41a-fbd192f368e1",
                    "wi3.sharepoint.com,d76580e2-d0c7-4af7-a545-3c2fb4c6ec23,8af0b1db-dc2c-4c41-a095-79bec35bd456",
                    "wi3.sharepoint.com,58ce2ee1-ffd1-4ed5-ae53-062a52aed104,8e063c0b-9e4b-4f72-9cfe-562dead72133",
                    "wi3.sharepoint.com,d9671191-fdd2-4104-aba3-f56112fffd74,d6045b5d-5dc7-4ae5-a9c0-e786c21606de",
                    "wi3.sharepoint.com,489d23a9-8f6d-4b9e-9fdf-0f9a6b0fc04d,ccb943f3-43e5-4d0f-b41a-fbd192f368e1",
                    "wi3.sharepoint.com,563a9c96-af9c-4d26-85ba-8cf9e7e2bde4,ccb943f3-43e5-4d0f-b41a-fbd192f368e1",
                    "wi3.sharepoint.com,f0021f9d-7332-4812-a3f3-ab410e06aabc,a59f0326-2de6-4360-8a55-1d0c89783f0e",
                    "wi3.sharepoint.com,ecdc38c6-24ca-4336-b661-c755025619fe,c394f607-f582-42e0-b859-61af215981ea",
                    "wi3.sharepoint.com,f4271222-fe3b-4379-9fc7-994559d55642,787f9249-8168-430f-9ac7-564ee0b8b3e5",
                    "wi3.sharepoint.com,22f9755f-6106-473c-b67a-39a4a235364c,cad904b1-213e-468b-a2eb-44926d5581c2",
                    "wi3.sharepoint.com,fa8f255e-926c-40af-b9d7-7fb16d0447e3,df2f0601-4214-4e1c-bf69-d7b2f714d85",
                    "wi3.sharepoint.com,ae77b56d-7128-4a29-a185-b01625476b0d,5f385d5c-0017-4e13-8463-c399bf3b26f7",
                    "wi3.sharepoint.com,a5c095dc-af5e-4645-80e7-12fe8dba9068,b445c790-280c-4c02-b258-3a39f0e268ba",
                    "wi3.sharepoint.com,a9bf53c5-fb96-48ee-9ad2-f29ea4252cff,0bb49602-528b-47d4-8d64-3947892413de",
                    "wi3.sharepoint.com,616f7cb9-f4b8-4904-b119-8ae56f593289,5f385d5c-0017-4e13-8463-c399bf3b26f7",
                    "wi3.sharepoint.com,729cb4c5-92d2-4645-80e9-dcf9ea517ddd,5f385d5c-0017-4e13-8463-c399bf3b26f7",
                    "wi3.sharepoint.com,d4377178-2910-484f-a942-0f19160dddf4,5fb4d46e-8a1c-4507-952b-3a5609f13760"
                    ]
def get_all_Sites():
    url = "https://graph.microsoft.com/v1.0/sites?search=*"
    sites = []

    while url:
        response = safe_get(url, get_headers())
        data = response.json()
        for site in data.get("value", []):
            sites.append(site["id"])

        url = data.get("@odata.nextLink")
    return sites

def check_if_sites_exists(sites_list):
    all_sites = get_all_Sites()
    exisiting_sites = set(all_sites)

    return [site_id for site_id in sites_list if site_id in exisiting_sites]



def process_all_sites(sites_list):
    for site in sites_list:
        print("porcess site")
        response = safe_get(f"https://graph.microsoft.com/v1.0/sites/{site}/drives", get_headers())
        drives = response.json()["value"]
        for drive in drives:
            if drive.get("name") == "Dokumente":
                print("process root")
                process_root(site, drive["id"])
                break
    return

def process_root(siteId, driveId):
    url = f"https://graph.microsoft.com/v1.0/sites/{siteId}/drives/{driveId}/root/children"
    while url:
        response = safe_get(url, get_headers())
        data = response.json()
        for item in data.get("value", []):
            if "folder" in item:
                print("Ordner:", item["name"])
                process_folder(siteId, driveId, item["id"])

            elif "file" in item:
                print("Datei:", item["name"])
                process_file(siteId, driveId, item["id"])
        url = data.get("@odata.nextLink")
    
def process_folder(siteId, driveId, folderId):
    url = f"https://graph.microsoft.com/v1.0/sites/{siteId}/drives/{driveId}/items/{folderId}/children"

    while url:
        response = safe_get(url, get_headers())
        data = response.json()
        for item in data.get("value", []):
            if "folder" in item:
                print("process folder", item["name"])
                process_folder(siteId, driveId, item["id"])
  
            elif "file" in item:
                print("process file", item["name"])
                process_file(siteId, driveId, item["id"])
        url = data.get("@odata.nextLink")


def process_file(siteId, driveId, item):
    response = safe_get(f"https://graph.microsoft.com/v1.0/sites/{siteId}/drives/{driveId}/items/{item}", get_headers())
    metaData = response.json()

    existing = cursor.execute("""SELECT modified_at FROM documents WHERE id=?""", (metaData["id"],)).fetchone()
    if existing and existing[0] == metaData["lastModifiedDateTime"]:
        return

    supported_types = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }

    if metaData["file"]["mimeType"] not in supported_types:
        return
    
    response = safe_get(f"https://graph.microsoft.com/v1.0/sites/{siteId}/drives/{driveId}/items/{item}/content", get_headers())
    text = ""
    try:
        match metaData["file"]["mimeType"]:
            case "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                document = Document(BytesIO(response.content))
                paragraphs = []
                for paragraph in document.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text)
                text = "\n".join(paragraphs)

            case "application/pdf":
                document = PdfReader(BytesIO(response.content))
                text_parts = []
                for page in document.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                text = "\n".join(text_parts)
        
            case "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                document = load_workbook(BytesIO(response.content))
                text_parts = []
                for sheet in document.worksheets:
                    text_parts.append(f"Sheet: {sheet.title}")
                    rows = sheet.iter_rows(values_only=True)
                    col_headers = next(rows, None)

                    if not col_headers:
                        continue
                    for row in rows:
                        values = []

                        for header, value in zip(col_headers,row):
                            if value is not None:
                                values.append(f"{header}: {value}")
                        if values:
                            text_parts.append(" | ".join(values))
                    text_parts.append("")
                text = "\n".join(text_parts)

            case "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                document = Presentation(BytesIO(response.content))
                text_parts = []
                for slide_number, slide in enumerate(document.slides, start=1):
                    text_parts.append(f"Slide: {slide_number}")

                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            shape_text = shape.text.strip()

                            if shape_text:
                                text_parts.append(shape_text)
                    text_parts.append("")
                text = "\n".join(text_parts)
            case _:
                return
    except Exception as e:
        print(f"Fehler beim Verarbeiten von {metaData['name']}: {e}")
        return
        
    cursor.execute("""
        INSERT INTO documents(
            id,
            site_id,
            drive_id,
            name,
            mime_type,
            size,
            web_url,
            created_at,
            modified_at,
            parent_id,
            text
        )
        VALUES(?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)

        ON CONFLICT(id) DO UPDATE SET
        site_id=excluded.site_id,
        drive_id=excluded.drive_id,
        name=excluded.name,
        mime_type=excluded.mime_type,
        size=excluded.size,
        web_url=excluded.web_url,
        created_at=excluded.created_at,
        modified_at=excluded.modified_at,
        parent_id=excluded.parent_id,
        text=excluded.text
    """, (
        metaData["id"],
        metaData["parentReference"]["siteId"],
        metaData["parentReference"]["driveId"],
        metaData["name"],
        metaData["file"]["mimeType"],
        metaData["size"],
        metaData["webUrl"],
        metaData["createdDateTime"],
        metaData["lastModifiedDateTime"],
        metaData["parentReference"]["id"],
        text
    ))
    connection.commit()
    return


process_all_sites(check_if_sites_exists(sites_to_process))
chunking.process_documents()

connection.close()